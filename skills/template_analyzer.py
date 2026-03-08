# -*- coding: utf-8 -*-
"""
ArcMind — Template Analyzer
==============================
從 PPTX / XLSX 提取「設計 DNA」，存為 JSON 設計檔。
任何公司上傳模板 → 分析 → 之後生成文件都帶有公司風格。
"""
from __future__ import annotations

import json
import logging
import os
import shutil
from pathlib import Path
from typing import Any

logger = logging.getLogger("arcmind.template_analyzer")

TEMPLATES_DIR = Path(__file__).parent.parent / "data" / "templates"
TEMPLATES_DIR.mkdir(parents=True, exist_ok=True)


# ── PPT Analysis ─────────────────────────────────────────────────────────────

def analyze_pptx(file_path: str, template_name: str) -> dict:
    """Analyze a PPTX file and extract its design DNA."""
    from pptx import Presentation
    from pptx.util import Pt, Emu

    prs = Presentation(file_path)
    profile: dict[str, Any] = {
        "name": template_name,
        "type": "pptx",
        "source_file": os.path.basename(file_path),
        "dimensions": {
            "width_emu": prs.slide_width,
            "height_emu": prs.slide_height,
            "width_inches": round(prs.slide_width / 914400, 1),
            "height_inches": round(prs.slide_height / 914400, 1),
            "aspect": "16:9" if abs(prs.slide_width / prs.slide_height - 16 / 9) < 0.1 else "4:3",
        },
        "slide_count": len(prs.slides),
        "fonts": set(),
        "colors": set(),
        "font_sizes": set(),
        "layouts": [],
        "slides": [],
    }

    # Analyze layouts
    for layout in prs.slide_layouts:
        ph_info = []
        for ph in layout.placeholders:
            ph_info.append({
                "idx": ph.placeholder_format.idx,
                "type": str(ph.placeholder_format.type),
                "name": ph.name,
            })
        profile["layouts"].append({
            "name": layout.name,
            "placeholders": ph_info,
        })

    # Analyze each slide
    for slide_idx, slide in enumerate(prs.slides):
        slide_info = {
            "index": slide_idx + 1,
            "layout": slide.slide_layout.name,
            "elements": [],
        }

        for shape in slide.shapes:
            elem = {
                "name": shape.name,
                "type": str(shape.shape_type),
                "position": {
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                },
            }

            # Extract text formatting
            if shape.has_text_frame:
                text_preview = shape.text[:100].replace("\n", " | ")
                elem["text_preview"] = text_preview
                elem["paragraphs"] = []
                for para in shape.text_frame.paragraphs:
                    para_info = {"alignment": str(para.alignment)}
                    for run in para.runs:
                        if run.font.name:
                            profile["fonts"].add(run.font.name)
                        if run.font.size:
                            profile["font_sizes"].add(run.font.size)
                        try:
                            if run.font.color and run.font.color.type is not None:
                                rgb = run.font.color.rgb
                                if rgb:
                                    profile["colors"].add(str(rgb))
                        except (AttributeError, TypeError):
                            pass
                    if para.runs:
                        r = para.runs[0]
                        para_info["font"] = r.font.name
                        para_info["size_pt"] = round(r.font.size / 12700, 1) if r.font.size else None
                        para_info["bold"] = r.font.bold
                    elem["paragraphs"].append(para_info)

            # Table info
            if shape.has_table:
                table = shape.table
                rows = len(list(table.rows))
                cols = len(list(table.columns))
                header = [cell.text[:20] for cell in list(table.rows)[0].cells]
                elem["table"] = {"rows": rows, "cols": cols, "header": header}

            # Image info
            try:
                if hasattr(shape, "image"):
                    elem["has_image"] = True
                    elem["image_type"] = shape.image.content_type
            except Exception:
                pass

            slide_info["elements"].append(elem)

        profile["slides"].append(slide_info)

    # Convert sets to sorted lists for JSON serialization
    profile["fonts"] = sorted(profile["fonts"])
    profile["colors"] = sorted(profile["colors"])
    profile["font_sizes"] = sorted(profile["font_sizes"])

    # Save profile
    profile_dir = TEMPLATES_DIR / template_name
    profile_dir.mkdir(parents=True, exist_ok=True)

    with open(profile_dir / "design_profile.json", "w", encoding="utf-8") as f:
        json.dump(profile, f, ensure_ascii=False, indent=2)

    # Copy original template
    dest = profile_dir / os.path.basename(file_path)
    if not dest.exists():
        shutil.copy2(file_path, dest)

    logger.info("[TemplateAnalyzer] PPTX '%s' analyzed: %d slides, fonts=%s, colors=%s",
                template_name, len(prs.slides), profile["fonts"], profile["colors"])

    return {
        "template_name": template_name,
        "type": "pptx",
        "slides": len(prs.slides),
        "fonts": profile["fonts"],
        "colors": profile["colors"],
        "layouts": len(profile["layouts"]),
        "profile_path": str(profile_dir / "design_profile.json"),
    }


# ── Excel Analysis ───────────────────────────────────────────────────────────

def analyze_xlsx(file_path: str, template_name: str) -> dict:
    """Analyze an XLSX file and extract its design DNA."""
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter

    wb = load_workbook(file_path, data_only=True)
    profile: dict[str, Any] = {
        "name": template_name,
        "type": "xlsx",
        "source_file": os.path.basename(file_path),
        "sheets": [],
        "fonts": set(),
        "colors": set(),
        "font_sizes": set(),
    }

    for ws in wb.worksheets:
        sheet_info = {
            "name": ws.title,
            "dimensions": ws.dimensions,
            "max_row": ws.max_row,
            "max_col": ws.max_column,
            "column_widths": {},
            "header_style": None,
            "merge_ranges": [str(m) for m in ws.merged_cells.ranges],
        }

        # Column widths
        for col_idx in range(1, min(ws.max_column or 1, 20) + 1):
            letter = get_column_letter(col_idx)
            dim = ws.column_dimensions.get(letter)
            if dim and dim.width:
                sheet_info["column_widths"][letter] = dim.width

        # Analyze header row style
        if ws.max_row and ws.max_row > 0:
            header_cells = []
            for cell in ws[1]:
                cell_info = {"value": str(cell.value or "")[:30]}
                if cell.font:
                    if cell.font.name:
                        profile["fonts"].add(cell.font.name)
                    cell_info["font"] = cell.font.name
                    cell_info["bold"] = cell.font.bold
                    cell_info["size"] = cell.font.size
                    if cell.font.size:
                        profile["font_sizes"].add(cell.font.size)
                if cell.fill and cell.fill.fgColor and cell.fill.fgColor.rgb:
                    rgb = str(cell.fill.fgColor.rgb)
                    if rgb != "00000000":
                        profile["colors"].add(rgb)
                        cell_info["bg_color"] = rgb
                if cell.alignment:
                    cell_info["alignment"] = cell.alignment.horizontal
                header_cells.append(cell_info)
            sheet_info["header_style"] = header_cells

        # Sample data rows for style
        data_styles = []
        for row_idx in range(2, min(ws.max_row or 1, 5) + 1):
            for cell in ws[row_idx]:
                if cell.font and cell.font.name:
                    profile["fonts"].add(cell.font.name)
                if cell.fill and cell.fill.fgColor and cell.fill.fgColor.rgb:
                    rgb = str(cell.fill.fgColor.rgb)
                    if rgb != "00000000":
                        profile["colors"].add(rgb)

        sheet_info["sample_data_styles"] = data_styles
        profile["sheets"].append(sheet_info)

    # Convert sets
    profile["fonts"] = sorted(profile["fonts"])
    profile["colors"] = sorted(profile["colors"])
    profile["font_sizes"] = sorted(profile["font_sizes"])

    # Save
    profile_dir = TEMPLATES_DIR / template_name
    profile_dir.mkdir(parents=True, exist_ok=True)

    with open(profile_dir / "design_profile.json", "w", encoding="utf-8") as f:
        json.dump(profile, f, ensure_ascii=False, indent=2)

    dest = profile_dir / os.path.basename(file_path)
    if not dest.exists():
        shutil.copy2(file_path, dest)

    logger.info("[TemplateAnalyzer] XLSX '%s' analyzed: %d sheets, fonts=%s",
                template_name, len(wb.worksheets), profile["fonts"])

    return {
        "template_name": template_name,
        "type": "xlsx",
        "sheets": [s["name"] for s in profile["sheets"]],
        "fonts": profile["fonts"],
        "colors": profile["colors"],
        "profile_path": str(profile_dir / "design_profile.json"),
    }


# ── Utilities ────────────────────────────────────────────────────────────────

def list_templates() -> list[dict]:
    """List all analyzed templates."""
    templates = []
    for d in TEMPLATES_DIR.iterdir():
        if d.is_dir():
            profile_path = d / "design_profile.json"
            if profile_path.exists():
                with open(profile_path, "r", encoding="utf-8") as f:
                    p = json.load(f)
                templates.append({
                    "name": p["name"],
                    "type": p["type"],
                    "fonts": p.get("fonts", []),
                    "colors": p.get("colors", []),
                    "source": p.get("source_file", ""),
                })
    return templates


def get_profile(template_name: str) -> dict | None:
    """Load a template's design profile."""
    profile_path = TEMPLATES_DIR / template_name / "design_profile.json"
    if profile_path.exists():
        with open(profile_path, "r", encoding="utf-8") as f:
            return json.load(f)
    return None


def get_template_file(template_name: str) -> str | None:
    """Get the original template file path."""
    tpl_dir = TEMPLATES_DIR / template_name
    if not tpl_dir.exists():
        return None
    for f in tpl_dir.iterdir():
        if f.suffix in (".pptx", ".xlsx") and f.name != "design_profile.json":
            return str(f)
    return None
