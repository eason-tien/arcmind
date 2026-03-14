# -*- coding: utf-8 -*-
"""
ArcMind Skill: Document Generation Toolbox
=============================================
提供原子操作讓 Agent 自由組合，而不是固定流程。
Agent 用腦力決定排版、配色、內容，此 skill 只負責執行。

兩種模式：
  1. 公司模板 — template=名稱 → 基於已學習的模板
  2. 自由設計 — 不指定 template → Agent 自由發揮

Actions 分為：模板管理 / PPT 操作 / Excel 操作
"""
from __future__ import annotations

import json
import logging
import os
from datetime import datetime
from pathlib import Path
from typing import Any

logger = logging.getLogger("arcmind.skills.document")

OUTPUT_DIR = Path(__file__).parent.parent / "data" / "output"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# In-memory workspace: active documents being worked on
_workspace: dict[str, Any] = {}


# ═══════════════════════════════════════════════════════════════════════════════
#  Template Management
# ═══════════════════════════════════════════════════════════════════════════════

def _analyze_template(inputs: dict) -> dict:
    """Analyze a company template file."""
    file_path = inputs.get("path", "")
    name = inputs.get("name", "")
    if not file_path or not name:
        return {"error": "Required: path (file), name (template name)"}
    if not os.path.exists(file_path):
        return {"error": f"File not found: {file_path}"}

    from skills.template_analyzer import analyze_pptx, analyze_xlsx
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pptx":
        return analyze_pptx(file_path, name)
    elif ext in (".xlsx", ".xls"):
        return analyze_xlsx(file_path, name)
    return {"error": f"Unsupported: {ext}. Use .pptx or .xlsx"}


def _list_templates(inputs: dict) -> dict:
    from skills.template_analyzer import list_templates
    templates = list_templates()
    return {"templates": templates, "count": len(templates)}


def _get_template_info(inputs: dict) -> dict:
    """Get detailed design profile of a template."""
    from skills.template_analyzer import get_profile
    name = inputs.get("name", "")
    profile = get_profile(name)
    if not profile:
        return {"error": f"Template '{name}' not found"}
    return profile


# ═══════════════════════════════════════════════════════════════════════════════
#  PPT Operations — 原子操作，Agent 決定怎麼用
# ═══════════════════════════════════════════════════════════════════════════════

def _ppt_new(inputs: dict) -> dict:
    """Create a new PPT workspace (from template or blank)."""
    from pptx import Presentation
    from pptx.util import Inches

    doc_id = inputs.get("doc_id", f"ppt_{datetime.now().strftime('%H%M%S')}")
    template_name = inputs.get("template", "")

    if template_name:
        from skills.template_analyzer import get_template_file
        tpl = get_template_file(template_name)
        if not tpl:
            return {"error": f"Template '{template_name}' not found"}
        prs = Presentation(tpl)
        mode = "company_template"
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        mode = "free_design"

    _workspace[doc_id] = {"type": "pptx", "prs": prs, "mode": mode}
    return {
        "doc_id": doc_id,
        "mode": mode,
        "slides": len(prs.slides),
        "layouts": [l.name for l in prs.slide_layouts],
    }


def _ppt_add_slide(inputs: dict) -> dict:
    """Add a slide. Agent decides layout."""
    doc_id = inputs.get("doc_id", "")
    ws = _workspace.get(doc_id)
    if not ws or ws["type"] != "pptx":
        return {"error": f"No active PPT: {doc_id}"}

    prs = ws["prs"]
    layout_name = inputs.get("layout", "")
    layout = None
    for l in prs.slide_layouts:
        if layout_name and (layout_name.lower() in l.name.lower()):
            layout = l
            break
    if not layout:
        # Default to blank or first available
        for l in prs.slide_layouts:
            if "blank" in l.name.lower() or "空白" in l.name:
                layout = l
                break
        if not layout:
            layout = prs.slide_layouts[0]

    slide = prs.slides.add_slide(layout)
    slide_idx = len(prs.slides)
    return {"slide_index": slide_idx, "layout": layout.name}


def _ppt_add_text(inputs: dict) -> dict:
    """Add text box to a slide. Agent controls position, font, color, size."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    doc_id = inputs.get("doc_id", "")
    ws = _workspace.get(doc_id)
    if not ws:
        return {"error": f"No active PPT: {doc_id}"}

    prs = ws["prs"]
    slide_idx = inputs.get("slide", len(prs.slides)) - 1
    if slide_idx < 0 or slide_idx >= len(prs.slides):
        return {"error": f"Invalid slide: {slide_idx + 1}"}
    slide = prs.slides[slide_idx]

    # Position (in inches, default to reasonable values)
    left = Inches(inputs.get("left", 0.5))
    top = Inches(inputs.get("top", 0.5))
    width = Inches(inputs.get("width", 9))
    height = Inches(inputs.get("height", 1))

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    text = inputs.get("text", "")
    font_name = inputs.get("font", "Arial")
    font_size = inputs.get("size", 14)
    bold = inputs.get("bold", False)
    color = inputs.get("color", "")

    lines = text.split("\n") if text else [""]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.bold = bold
        if color:
            try:
                p.font.color.rgb = RGBColor.from_string(color)
            except Exception:
                pass
        if inputs.get("line_spacing"):
            p.space_after = Pt(inputs["line_spacing"])

    return {"slide": slide_idx + 1, "shape": txBox.name}


def _ppt_add_table(inputs: dict) -> dict:
    """Add table to a slide. Agent provides data."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    doc_id = inputs.get("doc_id", "")
    ws = _workspace.get(doc_id)
    if not ws:
        return {"error": f"No active PPT: {doc_id}"}

    prs = ws["prs"]
    slide_idx = inputs.get("slide", len(prs.slides)) - 1
    slide = prs.slides[slide_idx]

    rows_data = inputs.get("rows", [])
    if not rows_data:
        return {"error": "rows required: [[col1, col2, ...], ...]"}

    n_rows = len(rows_data)
    n_cols = max(len(r) for r in rows_data)

    left = Inches(inputs.get("left", 0.5))
    top = Inches(inputs.get("top", 1.5))
    width = Inches(inputs.get("width", 11))
    height = Inches(inputs.get("height", 0.4 * n_rows))

    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    font_name = inputs.get("font", "Arial")
    header_color = inputs.get("header_color", "")

    for r_idx, row in enumerate(rows_data):
        for c_idx, val in enumerate(row):
            if c_idx >= n_cols:
                continue
            cell = tbl.cell(r_idx, c_idx)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.name = font_name
                p.font.size = Pt(inputs.get("font_size", 11))
                if r_idx == 0:
                    p.font.bold = True
                    if header_color:
                        try:
                            p.font.color.rgb = RGBColor.from_string(header_color)
                        except Exception:
                            pass

    return {"slide": slide_idx + 1, "rows": n_rows, "cols": n_cols}


def _ppt_add_shape(inputs: dict) -> dict:
    """Add a shape (rectangle, circle, arrow, etc.) to a slide."""
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    doc_id = inputs.get("doc_id", "")
    ws = _workspace.get(doc_id)
    if not ws:
        return {"error": f"No active PPT: {doc_id}"}

    prs = ws["prs"]
    slide_idx = inputs.get("slide", len(prs.slides)) - 1
    slide = prs.slides[slide_idx]

    shape_map = {
        "rectangle": MSO_SHAPE.RECTANGLE,
        "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
        "oval": MSO_SHAPE.OVAL,
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
        "arrow_right": MSO_SHAPE.RIGHT_ARROW,
        "arrow_left": MSO_SHAPE.LEFT_ARROW,
        "pentagon": MSO_SHAPE.PENTAGON,
        "chevron": MSO_SHAPE.CHEVRON,
    }

    shape_type = shape_map.get(inputs.get("shape", "rectangle"), MSO_SHAPE.RECTANGLE)

    shape = slide.shapes.add_shape(
        shape_type,
        Inches(inputs.get("left", 0.5)),
        Inches(inputs.get("top", 0.5)),
        Inches(inputs.get("width", 2)),
        Inches(inputs.get("height", 1)),
    )

    fill_color = inputs.get("fill_color", "")
    if fill_color:
        shape.fill.solid()
        try:
            shape.fill.fore_color.rgb = RGBColor.from_string(fill_color)
        except Exception:
            pass

    if inputs.get("no_border", False):
        shape.line.fill.background()

    text = inputs.get("text", "")
    if text:
        from pptx.util import Pt
        shape.text = text
        for p in shape.text_frame.paragraphs:
            p.font.name = inputs.get("font", "Arial")
            p.font.size = Pt(inputs.get("font_size", 12))
            if inputs.get("text_color"):
                try:
                    p.font.color.rgb = RGBColor.from_string(inputs["text_color"])
                except Exception:
                    pass

    return {"slide": slide_idx + 1, "shape": shape.name}


def _ppt_add_image(inputs: dict) -> dict:
    """Add an image to a slide."""
    from pptx.util import Inches

    doc_id = inputs.get("doc_id", "")
    ws = _workspace.get(doc_id)
    if not ws:
        return {"error": f"No active PPT: {doc_id}"}

    prs = ws["prs"]
    slide_idx = inputs.get("slide", len(prs.slides)) - 1
    slide = prs.slides[slide_idx]

    image_path = inputs.get("image_path", "")
    if not image_path or not os.path.exists(image_path):
        return {"error": f"Image not found: {image_path}"}

    slide.shapes.add_picture(
        image_path,
        Inches(inputs.get("left", 0.5)),
        Inches(inputs.get("top", 1)),
        Inches(inputs.get("width", 4)),
        Inches(inputs.get("height", 3)),
    )
    return {"slide": slide_idx + 1, "image": image_path}


def _ppt_save(inputs: dict) -> dict:
    """Save the PPT to file."""
    doc_id = inputs.get("doc_id", "")
    ws = _workspace.get(doc_id)
    if not ws:
        return {"error": f"No active PPT: {doc_id}"}

    prs = ws["prs"]
    title = inputs.get("title", doc_id)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = inputs.get("filename", f"{title.replace(' ', '_')}_{timestamp}.pptx")
    output_path = OUTPUT_DIR / filename
    prs.save(str(output_path))

    # Clean up workspace
    if inputs.get("close", True):
        del _workspace[doc_id]

    return {"file": str(output_path), "slides": len(prs.slides)}


# ═══════════════════════════════════════════════════════════════════════════════
#  Excel Operations — 同樣原子操作
# ═══════════════════════════════════════════════════════════════════════════════

def _excel_new(inputs: dict) -> dict:
    """Create new Excel workspace."""
    from openpyxl import Workbook
    doc_id = inputs.get("doc_id", f"xl_{datetime.now().strftime('%H%M%S')}")
    wb = Workbook()
    _workspace[doc_id] = {"type": "xlsx", "wb": wb}
    return {"doc_id": doc_id, "sheets": [wb.active.title]}


def _excel_write(inputs: dict) -> dict:
    """Write data to cells. Agent controls everything."""
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    doc_id = inputs.get("doc_id", "")
    ws_data = _workspace.get(doc_id)
    if not ws_data:
        return {"error": f"No active Excel: {doc_id}"}

    wb = ws_data["wb"]
    sheet_name = inputs.get("sheet", wb.active.title)
    if sheet_name not in wb.sheetnames:
        wb.create_sheet(sheet_name)
    ws = wb[sheet_name]

    # Write headers
    headers = inputs.get("headers", [])
    start_row = inputs.get("start_row", 1)
    header_font = inputs.get("header_font", "Arial")
    header_color = inputs.get("header_bg", "0070C0")
    font_name = inputs.get("font", "Arial")

    if headers:
        h_font = Font(name=header_font, bold=True, size=12, color="FFFFFF")
        h_fill = PatternFill(start_color=header_color, end_color=header_color, fill_type="solid")
        for col_idx, h in enumerate(headers, 1):
            cell = ws.cell(row=start_row, column=col_idx, value=h)
            cell.font = h_font
            cell.fill = h_fill
            cell.alignment = Alignment(horizontal="center", vertical="center")

    # Write data rows
    data = inputs.get("data", [])
    data_start = start_row + (1 if headers else 0)
    alt_fill = PatternFill(start_color="F2F2F2", end_color="F2F2F2", fill_type="solid")

    for r_idx, row in enumerate(data):
        for c_idx, val in enumerate(row, 1):
            cell = ws.cell(row=data_start + r_idx, column=c_idx, value=val)
            cell.font = Font(name=font_name, size=11)
            cell.alignment = Alignment(horizontal="center", vertical="center")
            if r_idx % 2 == 1:
                cell.fill = alt_fill

    # Auto-width
    from openpyxl.utils import get_column_letter
    all_cols = max(len(headers), max((len(r) for r in data), default=0)) if headers or data else 0
    for col_idx in range(1, all_cols + 1):
        max_len = 8
        for row_idx in range(start_row, data_start + len(data)):
            val = ws.cell(row=row_idx, column=col_idx).value
            if val:
                max_len = max(max_len, len(str(val)))
        ws.column_dimensions[get_column_letter(col_idx)].width = min(max_len + 4, 40)

    return {"sheet": sheet_name, "rows_written": len(data), "headers": len(headers)}


def _excel_format(inputs: dict) -> dict:
    """Apply formatting to a range. Agent specifies exactly what."""
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    doc_id = inputs.get("doc_id", "")
    ws_data = _workspace.get(doc_id)
    if not ws_data:
        return {"error": f"No active Excel: {doc_id}"}

    wb = ws_data["wb"]
    ws = wb[inputs.get("sheet", wb.active.title)]

    cell_range = inputs.get("range", "A1")  # e.g. "A1:D5"
    cells = ws[cell_range]

    # Flatten if single cell
    if not isinstance(cells, tuple):
        cells = ((cells,),)
    elif cells and not isinstance(cells[0], tuple):
        cells = (cells,)

    count = 0
    for row in cells:
        for cell in row:
            if inputs.get("font"):
                cell.font = Font(
                    name=inputs.get("font", "Arial"),
                    size=inputs.get("size", 11),
                    bold=inputs.get("bold", False),
                    color=inputs.get("color", "000000"),
                )
            if inputs.get("bg_color"):
                cell.fill = PatternFill(
                    start_color=inputs["bg_color"],
                    end_color=inputs["bg_color"],
                    fill_type="solid",
                )
            if inputs.get("merge_title"):
                cell.alignment = Alignment(horizontal="center", vertical="center")
            count += 1

    return {"cells_formatted": count}


def _excel_save(inputs: dict) -> dict:
    """Save Excel to file."""
    doc_id = inputs.get("doc_id", "")
    ws_data = _workspace.get(doc_id)
    if not ws_data:
        return {"error": f"No active Excel: {doc_id}"}

    wb = ws_data["wb"]
    title = inputs.get("title", doc_id)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = inputs.get("filename", f"{title.replace(' ', '_')}_{timestamp}.xlsx")
    output_path = OUTPUT_DIR / filename
    wb.save(str(output_path))

    if inputs.get("close", True):
        del _workspace[doc_id]

    return {"file": str(output_path), "sheets": wb.sheetnames}


# ═══════════════════════════════════════════════════════════════════════════════
#  Main Entry Point
# ═══════════════════════════════════════════════════════════════════════════════

_ACTIONS = {
    # Template management
    "analyze_template": _analyze_template,
    "list_templates": _list_templates,
    "get_template_info": _get_template_info,
    # PPT atomic operations
    "ppt_new": _ppt_new,
    "ppt_add_slide": _ppt_add_slide,
    "ppt_add_text": _ppt_add_text,
    "ppt_add_table": _ppt_add_table,
    "ppt_add_shape": _ppt_add_shape,
    "ppt_add_image": _ppt_add_image,
    "ppt_save": _ppt_save,
    # Excel atomic operations
    "excel_new": _excel_new,
    "excel_write": _excel_write,
    "excel_format": _excel_format,
    "excel_save": _excel_save,
}


def run(inputs: dict) -> dict:
    """
    Document Toolbox — Agent 用原子操作自由組合。

    模板管理:
      analyze_template, list_templates, get_template_info

    PPT 操作 (先 ppt_new → 各種 add → ppt_save):
      ppt_new          — 建立新 PPT (template= 或 空白)
      ppt_add_slide    — 新增頁面 (指定 layout)
      ppt_add_text     — 加文字框 (位置/字體/大小/顏色 全由 Agent 決定)
      ppt_add_table    — 加表格
      ppt_add_shape    — 加形狀 (矩形/圓/箭頭 + 填色 + 文字)
      ppt_add_image    — 加圖片
      ppt_save         — 儲存

    Excel 操作 (先 excel_new → write/format → excel_save):
      excel_new        — 建立新 Excel
      excel_write      — 寫入資料 (標題列 + 資料)
      excel_format     — 格式化儲存格
      excel_save       — 儲存
    """
    action = inputs.get("action", "")
    handler = _ACTIONS.get(action)
    if not handler:
        return {
            "error": f"Unknown action: {action}",
            "available": list(_ACTIONS.keys()),
            "usage": "先 ppt_new/excel_new → 操作 → save",
        }
    try:
        return handler(inputs)
    except Exception as e:
        logger.error("[Document] %s failed: %s", action, e)
        return {"status": "error", "action": action, "error": str(e)}
