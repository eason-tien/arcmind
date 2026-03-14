# -*- coding: utf-8 -*-
# Copyright (c) 2024-2026 ArcMind Contributors
# Licensed under the MIT License. See LICENSE file in the project root.
"""
Skill: email_kb_skill
PST 郵件知識庫 — 讀取 .pst 郵件檔，分類建立知識庫供 Agent 搜尋學習

流程: .pst → readpst → mbox → Python mailbox 解析 → LLM 整理 → SQLite FTS5
附件: 自動提取 Word/Excel/PPT/PDF 文字內容，一起存入知識庫
LLM: 為每封郵件生成摘要、分類標籤、語言偵測（原始資料不動）
"""
from __future__ import annotations

import email
import email.header
import email.utils
import io
import json
import logging
import mailbox
import os
import re
import shutil
import sqlite3
import subprocess
import tempfile
import time
from datetime import datetime
from pathlib import Path
from typing import Any

logger = logging.getLogger("arcmind.skill.email_kb")

_DATA_DIR = Path(__file__).resolve().parent.parent / "data" / "email_kb"
_ATTACH_DIR = _DATA_DIR / "attachments"
_DB_PATH = _DATA_DIR / "emails.db"


# ── Database ────────────────────────────────────────────────────────────────

def _get_db() -> sqlite3.Connection:
    """Get or create the email knowledge base database."""
    _DATA_DIR.mkdir(parents=True, exist_ok=True)
    _ATTACH_DIR.mkdir(parents=True, exist_ok=True)
    db = sqlite3.connect(str(_DB_PATH), timeout=10)
    db.row_factory = sqlite3.Row
    db.execute("PRAGMA journal_mode=WAL")

    db.execute("""
        CREATE TABLE IF NOT EXISTS emails (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            message_id TEXT UNIQUE,
            subject TEXT,
            sender TEXT,
            sender_email TEXT,
            recipients TEXT,
            date TEXT,
            date_ts INTEGER,
            year INTEGER,
            month INTEGER,
            folder TEXT,
            body TEXT,
            attachments TEXT,
            attachment_content TEXT,
            ai_summary TEXT,
            ai_tags TEXT,
            ai_language TEXT,
            pst_source TEXT,
            imported_at TEXT
        )
    """)

    # Attachments table — each attachment as a separate record
    db.execute("""
        CREATE TABLE IF NOT EXISTS attachments (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            email_id INTEGER,
            filename TEXT,
            file_type TEXT,
            file_path TEXT,
            file_size INTEGER,
            content_text TEXT,
            FOREIGN KEY (email_id) REFERENCES emails(id)
        )
    """)

    # FTS5 — now includes attachment_content
    db.execute("""
        CREATE VIRTUAL TABLE IF NOT EXISTS emails_fts USING fts5(
            subject, sender, body, folder, attachment_content, ai_summary, ai_tags,
            content='emails',
            content_rowid='id'
        )
    """)

    # Triggers
    db.executescript("""
        CREATE TRIGGER IF NOT EXISTS emails_ai AFTER INSERT ON emails BEGIN
            INSERT INTO emails_fts(rowid, subject, sender, body, folder, attachment_content, ai_summary, ai_tags)
            VALUES (new.id, new.subject, new.sender, new.body, new.folder, new.attachment_content, new.ai_summary, new.ai_tags);
        END;
        CREATE TRIGGER IF NOT EXISTS emails_ad AFTER DELETE ON emails BEGIN
            INSERT INTO emails_fts(emails_fts, rowid, subject, sender, body, folder, attachment_content, ai_summary, ai_tags)
            VALUES('delete', old.id, old.subject, old.sender, old.body, old.folder, old.attachment_content, old.ai_summary, old.ai_tags);
        END;
    """)

    # Migration: add columns if missing (upgrade path)
    for col in ("attachment_content", "ai_summary", "ai_tags", "ai_language"):
        try:
            db.execute(f"SELECT {col} FROM emails LIMIT 1")
        except sqlite3.OperationalError:
            db.execute(f"ALTER TABLE emails ADD COLUMN {col} TEXT")
            logger.info("[email_kb] Migrated: added %s column", col)

    db.commit()
    return db


# ── Attachment Content Extraction ───────────────────────────────────────────

def _extract_docx(data: bytes) -> str:
    """Extract text from .docx file."""
    try:
        from docx import Document
        doc = Document(io.BytesIO(data))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))
        return "\n".join(paragraphs)
    except Exception as e:
        logger.debug("[email_kb] docx extract error: %s", e)
        return ""


def _extract_xlsx(data: bytes) -> str:
    """Extract text from .xlsx file."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        lines = []
        for sheet in wb.sheetnames[:10]:  # Max 10 sheets
            ws = wb[sheet]
            lines.append(f"[Sheet: {sheet}]")
            row_count = 0
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    lines.append(" | ".join(cells))
                row_count += 1
                if row_count > 200:  # Cap rows per sheet
                    lines.append("... (截斷)")
                    break
        wb.close()
        return "\n".join(lines)
    except Exception as e:
        logger.debug("[email_kb] xlsx extract error: %s", e)
        return ""


def _extract_pptx(data: bytes) -> str:
    """Extract text from .pptx file."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                lines.append(f"[Slide {i}]")
                lines.extend(slide_texts)
        return "\n".join(lines)
    except Exception as e:
        logger.debug("[email_kb] pptx extract error: %s", e)
        return ""


def _extract_pdf(data: bytes) -> str:
    """Extract text from .pdf file."""
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(io.BytesIO(data))
        texts = []
        for page in reader.pages[:50]:  # Max 50 pages
            text = page.extract_text()
            if text:
                texts.append(text)
        return "\n".join(texts)
    except Exception as e:
        logger.debug("[email_kb] pdf extract error: %s", e)
        return ""


def _extract_text_file(data: bytes) -> str:
    """Extract text from plain text / CSV files."""
    for enc in ("utf-8", "big5", "gb2312", "shift_jis", "latin-1"):
        try:
            return data.decode(enc)[:20000]
        except (UnicodeDecodeError, LookupError):
            continue
    return ""


# Map file extensions to extractors
_EXTRACTORS = {
    ".docx": _extract_docx,
    ".doc": None,  # Binary .doc not supported without antiword
    ".xlsx": _extract_xlsx,
    ".xls": None,   # Binary .xls not directly supported
    ".pptx": _extract_pptx,
    ".ppt": None,    # Binary .ppt not directly supported
    ".pdf": _extract_pdf,
    ".txt": _extract_text_file,
    ".csv": _extract_text_file,
    ".log": _extract_text_file,
    ".md": _extract_text_file,
    ".json": _extract_text_file,
    ".xml": _extract_text_file,
    ".html": _extract_text_file,
    ".htm": _extract_text_file,
}


def _extract_attachment_content(filename: str, data: bytes) -> str:
    """Extract text content from an attachment based on file extension."""
    ext = Path(filename).suffix.lower()
    extractor = _EXTRACTORS.get(ext)
    if extractor is None:
        return ""
    try:
        text = extractor(data)
        return text[:30000]  # Cap at 30K chars per attachment
    except Exception as e:
        logger.debug("[email_kb] Attachment extract error (%s): %s", filename, e)
        return ""


# ── Email Helpers ───────────────────────────────────────────────────────────

def _decode_header(raw: str | None) -> str:
    """Decode RFC2047 encoded email header."""
    if not raw:
        return ""
    parts = email.header.decode_header(raw)
    decoded = []
    for data, charset in parts:
        if isinstance(data, bytes):
            decoded.append(data.decode(charset or "utf-8", errors="replace"))
        else:
            decoded.append(str(data))
    return " ".join(decoded).strip()


def _extract_email_addr(raw: str | None) -> str:
    """Extract just the email address from a header value."""
    if not raw:
        return ""
    _, addr = email.utils.parseaddr(raw)
    return addr or raw


def _get_body(msg: email.message.Message) -> str:
    """Extract plain text body from email message."""
    body_parts: list[str] = []
    if msg.is_multipart():
        for part in msg.walk():
            ct = part.get_content_type()
            if ct == "text/plain":
                payload = part.get_payload(decode=True)
                if payload:
                    charset = part.get_content_charset() or "utf-8"
                    body_parts.append(payload.decode(charset, errors="replace"))
            elif ct == "text/html" and not body_parts:
                payload = part.get_payload(decode=True)
                if payload:
                    charset = part.get_content_charset() or "utf-8"
                    html = payload.decode(charset, errors="replace")
                    text = re.sub(r"<[^>]+>", " ", html)
                    text = re.sub(r"\s+", " ", text).strip()
                    body_parts.append(text)
    else:
        payload = msg.get_payload(decode=True)
        if payload:
            charset = msg.get_content_charset() or "utf-8"
            body_parts.append(payload.decode(charset, errors="replace"))
    return "\n".join(body_parts)[:10000]


def _process_attachments(msg: email.message.Message, email_id: int,
                         db: sqlite3.Connection) -> tuple[list[str], str]:
    """
    Extract, save, and parse all attachments from an email.
    Returns (attachment_names, combined_content_text).
    """
    att_names: list[str] = []
    all_content: list[str] = []

    if not msg.is_multipart():
        return att_names, ""

    for part in msg.walk():
        fn = part.get_filename()
        if not fn:
            continue

        filename = _decode_header(fn)
        att_names.append(filename)

        # Get raw attachment data
        payload = part.get_payload(decode=True)
        if not payload:
            continue

        # Save to disk
        safe_name = re.sub(r'[^\w\-.]', '_', filename)
        att_subdir = _ATTACH_DIR / str(email_id)
        att_subdir.mkdir(parents=True, exist_ok=True)
        file_path = att_subdir / safe_name

        try:
            file_path.write_bytes(payload)
        except Exception as e:
            logger.debug("[email_kb] Save attachment error: %s", e)
            continue

        # Extract text content
        content_text = _extract_attachment_content(filename, payload)

        # Store in DB
        try:
            db.execute("""
                INSERT INTO attachments (email_id, filename, file_type, file_path, file_size, content_text)
                VALUES (?, ?, ?, ?, ?, ?)
            """, (email_id, filename, Path(filename).suffix.lower(),
                  str(file_path), len(payload), content_text[:30000]))
        except Exception as e:
            logger.debug("[email_kb] DB attachment insert error: %s", e)

        if content_text:
            all_content.append(f"[附件: {filename}]\n{content_text}")

    combined = "\n\n".join(all_content)[:50000]
    return att_names, combined


def _parse_date(raw: str | None) -> tuple[str, int, int, int]:
    """Parse email date → (iso_str, timestamp, year, month)."""
    if not raw:
        return ("", 0, 0, 0)
    try:
        parsed = email.utils.parsedate_to_datetime(raw)
        return (parsed.isoformat(), int(parsed.timestamp()),
                parsed.year, parsed.month)
    except Exception:
        return (raw, 0, 0, 0)


# ── LLM Processing ─────────────────────────────────────────────────────────

def _llm_process_email(subject: str, sender: str, body: str,
                       att_content: str) -> dict:
    """
    Use LLM to generate structured knowledge from an email.
    Returns {ai_summary, ai_tags, ai_language}.
    NEVER modifies original data — only generates supplementary fields.
    """
    try:
        from runtime.model_router import model_router
    except Exception:
        return {"ai_summary": "", "ai_tags": "", "ai_language": ""}

    # Build concise input for LLM (limit to save tokens)
    body_preview = (body or "")[:2000]
    att_preview = (att_content or "")[:1000]

    prompt = f"""Analyze this email and return a JSON object with exactly 3 fields.
Do NOT add any text outside the JSON. Do NOT invent information not present.

Email:
- Subject: {subject}
- From: {sender}
- Body: {body_preview}
{f'- Attachment content: {att_preview}' if att_preview else ''}

Return JSON:
{{
  "summary": "One sentence summarizing what this email is about, in the same language as the email",
  "tags": "comma-separated category tags (e.g. 會議,月報,採購,技術文件,人事,通知,報價,客戶)",
  "language": "primary language code: zh, en, ja, or mixed"
}}"""

    try:
        resp = model_router.complete(
            prompt=prompt,
            task_type="classify",
            budget="low",
            max_tokens=300,
        )
        text = resp.content.strip()
        # Extract JSON from response
        if "{" in text and "}" in text:
            json_str = text[text.index("{"):text.rindex("}") + 1]
            data = json.loads(json_str)
            return {
                "ai_summary": str(data.get("summary", ""))[:500],
                "ai_tags": str(data.get("tags", ""))[:200],
                "ai_language": str(data.get("language", ""))[:10],
            }
    except Exception as e:
        logger.debug("[email_kb] LLM process error: %s", e)

    return {"ai_summary": "", "ai_tags": "", "ai_language": ""}


# ── Import ──────────────────────────────────────────────────────────────────

def _import_pst(inputs: dict) -> dict:
    """Import a .pst file into the knowledge base (with attachment extraction)."""
    pst_path = inputs.get("pst_path", "")
    if not pst_path:
        return {"success": False, "error": "pst_path 為必填"}
    if not Path(pst_path).exists():
        return {"success": False, "error": f"檔案不存在: {pst_path}"}

    logger.info("[email_kb] Importing PST: %s", pst_path)

    # Convert PST → mbox using readpst (with attachment extraction)
    tmp_dir = tempfile.mkdtemp(prefix="arcmind_pst_")
    try:
        cmd = ["readpst", "-r", "-o", tmp_dir, pst_path]
        result = subprocess.run(cmd, capture_output=True, timeout=1800)
        stderr_text = result.stderr.decode("utf-8", errors="replace")
        if result.returncode != 0:
            return {"success": False, "error": f"readpst 失敗: {stderr_text[:500]}"}

        mbox_files = [f for f in Path(tmp_dir).rglob("*")
                      if f.is_file() and f.stat().st_size > 0]

        if not mbox_files:
            return {"success": False, "error": "readpst 未產出任何檔案"}

        db = _get_db()
        total = 0
        skipped = 0
        errors = 0
        att_total = 0
        att_extracted = 0
        pst_name = Path(pst_path).name

        for mbox_path in mbox_files:
            rel = mbox_path.relative_to(tmp_dir)
            folder = str(rel.parent) if str(rel.parent) != "." else rel.stem

            try:
                mbox_obj = mailbox.mbox(str(mbox_path))
            except Exception:
                continue

            for msg in mbox_obj:
                try:
                    msg_id = msg.get("Message-ID", "")
                    subject = _decode_header(msg.get("Subject"))
                    sender = _decode_header(msg.get("From", ""))
                    sender_email = _extract_email_addr(msg.get("From", ""))
                    recipients = _decode_header(msg.get("To", ""))
                    date_raw = msg.get("Date", "")
                    date_str, date_ts, year, month = _parse_date(date_raw)
                    body = _get_body(msg)

                    if not msg_id:
                        msg_id = f"{pst_name}_{total}_{hash(subject + sender)}"

                    try:
                        cursor = db.execute("""
                            INSERT OR IGNORE INTO emails
                            (message_id, subject, sender, sender_email,
                             recipients, date, date_ts, year, month,
                             folder, body, attachments, attachment_content,
                             ai_summary, ai_tags, ai_language,
                             pst_source, imported_at)
                            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                        """, (msg_id, subject, sender, sender_email,
                              recipients, date_str, date_ts, year, month,
                              folder, body, "[]", "",
                              "", "", "",
                              pst_name, datetime.now().isoformat()))

                        if cursor.rowcount > 0:
                            email_id = cursor.lastrowid
                            # Process attachments
                            att_names, att_content = _process_attachments(
                                msg, email_id, db)

                            if att_names:
                                att_total += len(att_names)
                                if att_content:
                                    att_extracted += len(att_names)

                            # LLM processing
                            ai_data = _llm_process_email(
                                subject, sender, body, att_content)
                            ai_summary = ai_data["ai_summary"]
                            ai_tags = ai_data["ai_tags"]
                            ai_language = ai_data["ai_language"]

                            # Update with attachments + AI data
                            db.execute("""
                                UPDATE emails
                                SET attachments = ?, attachment_content = ?,
                                    ai_summary = ?, ai_tags = ?, ai_language = ?
                                WHERE id = ?
                            """, (json.dumps(att_names, ensure_ascii=False),
                                  att_content, ai_summary, ai_tags, ai_language,
                                  email_id))

                            # Update FTS (trigger fired with empty values)
                            db.execute("""
                                INSERT INTO emails_fts(emails_fts, rowid,
                                    subject, sender, body, folder,
                                    attachment_content, ai_summary, ai_tags)
                                VALUES('delete', ?, ?, ?, ?, ?, ?, ?, ?)
                            """, (email_id, subject, sender, body, folder,
                                  "", "", ""))
                            db.execute("""
                                INSERT INTO emails_fts(rowid,
                                    subject, sender, body, folder,
                                    attachment_content, ai_summary, ai_tags)
                                VALUES(?, ?, ?, ?, ?, ?, ?, ?)
                            """, (email_id, subject, sender, body, folder,
                                  att_content, ai_summary, ai_tags))

                            total += 1
                        else:
                            skipped += 1

                    except sqlite3.IntegrityError:
                        skipped += 1

                    if total > 0 and total % 200 == 0:
                        db.commit()
                        logger.info("[email_kb] Progress: %d emails, %d attachments",
                                    total, att_total)

                except Exception as e:
                    errors += 1
                    if errors <= 5:
                        logger.warning("[email_kb] Parse error: %s", e)

        db.commit()
        db.close()

        return {
            "success": True,
            "imported": total,
            "skipped_duplicates": skipped,
            "errors": errors,
            "attachments_found": att_total,
            "attachments_content_extracted": att_extracted,
            "pst_file": pst_name,
            "message": f"匯入完成！{total} 封郵件，{att_total} 個附件（{att_extracted} 個內容已提取）。",
        }

    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


# ── Search ──────────────────────────────────────────────────────────────────

def _search(inputs: dict) -> dict:
    """Full-text search across emails AND attachment content."""
    query = inputs.get("query", "")
    if not query:
        return {"success": False, "error": "query 為必填"}

    limit = int(inputs.get("limit", 10))
    db = _get_db()

    rows = db.execute("""
        SELECT e.id, e.subject, e.sender, e.date, e.folder, e.attachments,
               snippet(emails_fts, 2, '→', '←', '…', 30) as body_snippet,
               snippet(emails_fts, 4, '▶', '◀', '…', 30) as att_snippet
        FROM emails_fts
        JOIN emails e ON e.id = emails_fts.rowid
        WHERE emails_fts MATCH ?
        ORDER BY rank
        LIMIT ?
    """, (query, limit)).fetchall()

    results = []
    for r in rows:
        item = {
            "id": r["id"],
            "subject": r["subject"],
            "sender": r["sender"],
            "date": r["date"],
            "folder": r["folder"],
            "body_snippet": r["body_snippet"],
        }
        att_snippet = r["att_snippet"]
        if att_snippet and att_snippet.strip():
            item["attachment_snippet"] = att_snippet
        att_list = r["attachments"]
        if att_list and att_list != "[]":
            item["attachments"] = json.loads(att_list)
        results.append(item)

    db.close()
    return {
        "success": True,
        "query": query,
        "count": len(results),
        "results": results,
    }


# ── Stats ───────────────────────────────────────────────────────────────────

def _stats(inputs: dict) -> dict:
    """Get knowledge base statistics."""
    db = _get_db()

    total = db.execute("SELECT COUNT(*) as c FROM emails").fetchone()["c"]
    if total == 0:
        db.close()
        return {"success": True, "total": 0, "message": "知識庫為空，請先用 import_pst 匯入。"}

    att_total = db.execute("SELECT COUNT(*) as c FROM attachments").fetchone()["c"]
    att_with_content = db.execute(
        "SELECT COUNT(*) as c FROM attachments WHERE content_text IS NOT NULL AND content_text != ''"
    ).fetchone()["c"]

    by_year = db.execute(
        "SELECT year, COUNT(*) as c FROM emails WHERE year > 0 GROUP BY year ORDER BY year"
    ).fetchall()

    top_senders = db.execute(
        "SELECT sender_email, COUNT(*) as c FROM emails GROUP BY sender_email ORDER BY c DESC LIMIT 10"
    ).fetchall()

    by_folder = db.execute(
        "SELECT folder, COUNT(*) as c FROM emails GROUP BY folder ORDER BY c DESC LIMIT 10"
    ).fetchall()

    by_source = db.execute(
        "SELECT pst_source, COUNT(*) as c FROM emails GROUP BY pst_source"
    ).fetchall()

    att_by_type = db.execute(
        "SELECT file_type, COUNT(*) as c FROM attachments GROUP BY file_type ORDER BY c DESC LIMIT 10"
    ).fetchall()

    db.close()
    return {
        "success": True,
        "total_emails": total,
        "total_attachments": att_total,
        "attachments_with_content": att_with_content,
        "by_year": [{"year": r["year"], "count": r["c"]} for r in by_year],
        "top_senders": [{"email": r["sender_email"], "count": r["c"]} for r in top_senders],
        "by_folder": [{"folder": r["folder"], "count": r["c"]} for r in by_folder],
        "by_source": [{"source": r["pst_source"], "count": r["c"]} for r in by_source],
        "attachment_types": [{"type": r["file_type"], "count": r["c"]} for r in att_by_type],
    }


# ── Browse ──────────────────────────────────────────────────────────────────

def _browse(inputs: dict) -> dict:
    """Browse emails by category."""
    by = inputs.get("by", "folder")
    value = inputs.get("value", "")
    limit = int(inputs.get("limit", 20))
    offset = int(inputs.get("offset", 0))
    db = _get_db()

    if by == "folder":
        if value:
            rows = db.execute(
                "SELECT id, subject, sender, date, folder, attachments FROM emails WHERE folder = ? ORDER BY date_ts DESC LIMIT ? OFFSET ?",
                (value, limit, offset)).fetchall()
        else:
            rows = db.execute("SELECT folder, COUNT(*) as c FROM emails GROUP BY folder ORDER BY c DESC").fetchall()
            db.close()
            return {"success": True, "folders": [{"name": r["folder"], "count": r["c"]} for r in rows]}
    elif by == "sender":
        if value:
            rows = db.execute(
                "SELECT id, subject, sender, date, folder, attachments FROM emails WHERE sender_email LIKE ? ORDER BY date_ts DESC LIMIT ? OFFSET ?",
                (f"%{value}%", limit, offset)).fetchall()
        else:
            rows = db.execute("SELECT sender_email, COUNT(*) as c FROM emails GROUP BY sender_email ORDER BY c DESC LIMIT 30").fetchall()
            db.close()
            return {"success": True, "senders": [{"email": r["sender_email"], "count": r["c"]} for r in rows]}
    elif by == "year":
        if value:
            rows = db.execute(
                "SELECT id, subject, sender, date, folder, attachments FROM emails WHERE year = ? ORDER BY date_ts DESC LIMIT ? OFFSET ?",
                (int(value), limit, offset)).fetchall()
        else:
            rows = db.execute("SELECT year, COUNT(*) as c FROM emails WHERE year > 0 GROUP BY year ORDER BY year").fetchall()
            db.close()
            return {"success": True, "years": [{"year": r["year"], "count": r["c"]} for r in rows]}
    else:
        db.close()
        return {"success": False, "error": f"不支援的分類: {by}，可用 folder/sender/year"}

    db.close()
    return {"success": True, "by": by, "value": value, "emails": [dict(r) for r in rows], "count": len(rows)}


# ── Export ──────────────────────────────────────────────────────────────────

def _export(inputs: dict) -> dict:
    """Export emails (with attachment content) as Markdown knowledge documents."""
    query = inputs.get("query", "")
    by = inputs.get("by", "")
    value = inputs.get("value", "")
    limit = int(inputs.get("limit", 50))
    db = _get_db()

    if query:
        rows = db.execute("""
            SELECT e.* FROM emails_fts
            JOIN emails e ON e.id = emails_fts.rowid
            WHERE emails_fts MATCH ? ORDER BY rank LIMIT ?
        """, (query, limit)).fetchall()
    elif by and value:
        col_map = {"folder": "folder", "sender": "sender_email", "year": "year"}
        col = col_map.get(by, "folder")
        rows = db.execute(
            f"SELECT * FROM emails WHERE {col} = ? ORDER BY date_ts DESC LIMIT ?",
            (value, limit)).fetchall()
    else:
        db.close()
        return {"success": False, "error": "需指定 query 或 by+value"}

    if not rows:
        db.close()
        return {"success": True, "message": "無符合條件的郵件。", "count": 0}

    lines = [f"# 郵件知識庫匯出", "", f"查詢: {query or f'{by}={value}'}", f"共 {len(rows)} 封", ""]
    for r in rows:
        lines.append(f"## {r['subject'] or '(無主題)'}")
        lines.append(f"- **寄件者**: {r['sender']}")
        lines.append(f"- **日期**: {r['date']}")
        lines.append(f"- **資料夾**: {r['folder']}")
        body_preview = (r['body'] or '')[:500]
        if body_preview:
            lines.append(f"\n{body_preview}\n")
        att_content = r.get('attachment_content') or ''
        if att_content:
            lines.append(f"### 附件內容\n{att_content[:1000]}\n")
        lines.append("---\n")

    output_path = _DATA_DIR / f"export_{int(time.time())}.md"
    output_path.write_text("\n".join(lines), encoding="utf-8")
    db.close()
    return {"success": True, "path": str(output_path), "count": len(rows),
            "message": f"已匯出 {len(rows)} 封郵件到 {output_path.name}"}


# ── Read ────────────────────────────────────────────────────────────────────

def _read_email(inputs: dict) -> dict:
    """Read a single email by ID, including attachment content."""
    email_id = inputs.get("id", 0)
    if not email_id:
        return {"success": False, "error": "id 為必填"}
    db = _get_db()
    row = db.execute("SELECT * FROM emails WHERE id = ?", (email_id,)).fetchone()
    if not row:
        db.close()
        return {"success": False, "error": f"找不到 ID={email_id} 的郵件"}

    atts = db.execute(
        "SELECT filename, file_type, file_path, file_size, content_text FROM attachments WHERE email_id = ?",
        (email_id,)).fetchall()
    db.close()
    return {
        "success": True,
        "email": dict(row),
        "attachments_detail": [dict(a) for a in atts],
    }


# ── Entry Point ─────────────────────────────────────────────────────────────

def run(inputs: dict) -> dict:
    """
    Email Knowledge Base skill entry point.

    inputs:
      action: import_pst | search | stats | browse | export | read
      pst_path: PST 檔案路徑 (import_pst 用)
      query: 搜尋關鍵字 (search/export)
      by: 分類方式 folder|sender|year (browse/export)
      value: 分類值 (browse/export)
      id: 郵件 ID (read)
      limit: 結果數量上限

    搜尋範圍包含: 郵件主旨、寄件者、內文、資料夾、附件內容（Word/Excel/PPT/PDF）
    """
    action = inputs.get("action", "stats")
    handlers = {
        "import_pst": _import_pst,
        "search": _search,
        "stats": _stats,
        "browse": _browse,
        "export": _export,
        "read": _read_email,
    }
    handler = handlers.get(action)
    if not handler:
        return {"success": False, "error": f"未知 action: {action}",
                "available_actions": list(handlers.keys())}
    try:
        return handler(inputs)
    except Exception as e:
        logger.error("[email_kb] %s failed: %s", action, e)
        return {"success": False, "error": str(e), "action": action}
